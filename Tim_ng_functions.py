import itertools
import os
import json
from PyPDF2 import PdfMerger #need to write for pypdf
from pypdf import PdfWriter
from pypdf import PdfReader
import pandas as pd
import csv
import re
from pptx import Presentation
from pptx.util import Inches,Pt
from PIL import Image
import io
import fitz
import sys
import numpy as np
import matplotlib.pyplot as plt
import matplotlib.colors as colors
import matplotlib.colors as mcolors
from matplotlib import rcParams
import seaborn as sns
from Na12HMMModel_TF import *

def combine_pdfs(folder_path, out_sfx): #input folder pather where pdfs are stored, out_sfx = output suffix
    file = open('./JUPYTERmutant_list.txt','r') #put mutant names in here
    Lines = file.readlines() #read each mut name
    #os.chdir(folder_path)
    for line in Lines:
        mutTXT = line.strip() #remove return after each mut name
        print(mutTXT)
        mutfolderpath = os.path.join(folder_path,mutTXT)
        
        os.chdir(mutfolderpath)
        #
        #if filename.endswith(".pdf"):
            #merger.append(os.path.join(mutfolderpath, filename))
        merger = PdfMerger()
        #merger = PdfWriter() #for use with pypdf library *UNTESTED*

        # pdfs = [f'{mutTXT}__0.5_75.pdf',f'{mutTXT}__0.5_100.pdf',
        #         'test__fi.pdf','test_0.5_wtvmut.pdf','test_1_wtvmut.pdf']
        pdfs = ['test_0.5_wtvmut.pdf','test_1_wtvmut.pdf']
        
        for pdf in pdfs:
            merger.append(pdf)
        
        #merger.write(folder_path+"_combined100323.pdf")
        merger.write(folder_path+mutTXT+out_sfx+'.pdf')
    merger.close()

def combine_and_sort_ef_csvs(root_path, out_sfx):
  combined_ef = f'{root_path}combined_ef_{out_sfx}.csv'
  # Open the output file in write mode
  with open(combined_ef, "w", newline="") as outfile:
    writer = csv.writer(outfile)

    data_by_filename = {}

    # Loop through all folders in the current directory
    for folder in os.listdir(root_path):
      # Check if it's a directory and not a file
      if os.path.isdir(os.path.join(root_path, folder)):
        # Access files within the folder using another loop
        for filename in os.listdir(os.path.join(root_path, folder)):
          if filename.endswith(".csv"):
            # Extract mutation number and index from filename
            match = re.search(r"mut(\d+)_(\d+)_", filename)
            if match:
              mutation_number, index = match.groups()

              # Open and process the CSV file
              with open(os.path.join(root_path, folder, filename), "r", newline="") as infile:
                reader = csv.reader(infile)

                # Skip the header
                next(reader)

                # Store data row with filename information as key
                for row in reader:
                  data_by_filename[f"{mutation_number}_{index}_{filename}"] = row

    # Sort data based on filename key (automatically sorts based on structure)
    sorted_data = sorted(data_by_filename.items(), key=lambda item: item[0])

    # Write sorted data to the output file
    for key, row in sorted_data:
      writer.writerow(row)

  print(f"Successfully combined and sorted data into '{combined_ef}'!")
  return

def make_ppt_from_pdf2(pdf_path, output_ppt_path):
  prs = Presentation()
  blank_slide_layout = prs.slide_layouts[6]
  slide = prs.slides.add_slide(blank_slide_layout)
  left = top = width = height = Inches(0.5)
  
  
  # file = open("/global/homes/t/tfenton/Neuron_general-2/JUPYTERmutant_list.txt", "r")
  # lines = file.readlines()
  # for line in lines: #mut1_1,mut1_2 etc.
  #   mut_txt = line.strip()
  #   print(mut_txt)
    
  for folder in sorted(os.listdir(pdf_path)): #Sort the listed directories so output pptx is somewhat in order 
      if os.path.isdir(os.path.join(pdf_path, folder)):# Check if it's a directory
        slide = prs.slides.add_slide(blank_slide_layout) #Add slide for each folder. Each slide will have multiple plots on it
        txBox = slide.shapes.add_textbox(left,top,width,height) #Add text box for title
        tf = txBox.text_frame
        page_count = 0
        for filename in sorted(os.listdir(os.path.join(pdf_path, folder))):# Access files within the folder, and sort them to appear in sam order on slide
          if filename.endswith(".pdf"):
  
            file = os.path.join(pdf_path,folder, filename) #pdf file path
            
            tf.text = folder
            print(folder)
            
            doc = fitz.open(file)
            for page in doc:  # iterate through the pages
              pix = page.get_pixmap(dpi=150)  # render page to an image
              pix.save(f"{pdf_path}/{folder}-{filename}.png")  # store image as a PNG
              
              
              slide.shapes.add_picture(f"{pdf_path}/{folder}-{filename}.png",left=Inches(page_count*2.5),top=Inches(1), width=Inches(2.5))
              print(page_count)
          page_count+=1
          
  
  prs.save(output_ppt_path)
  print(f"Successfully converted PDFs to PowerPoint presentation: {output_ppt_path}")
  return

def make_ppt_from_pdf3(pdf_path, output_ppt_path):
  prs = Presentation()
  blank_slide_layout = prs.slide_layouts[6]
  slide = prs.slides.add_slide(blank_slide_layout)
  left = top = width = height = Inches(0.5)
  
  
  # file = open("/global/homes/t/tfenton/Neuron_general-2/JUPYTERmutant_list.txt", "r")
  # lines = file.readlines()
  # for line in lines: #mut1_1,mut1_2 etc.
  #   mut_txt = line.strip()
  #   print(mut_txt)
    
  for folder in sorted(os.listdir(pdf_path)): #Sort the listed directories so output pptx is somewhat in order 
      if os.path.isdir(os.path.join(pdf_path, folder)):# Check if it's a directory
        slide = prs.slides.add_slide(blank_slide_layout) #Add slide for each folder. Each slide will have multiple plots on it
        txBox = slide.shapes.add_textbox(left,top,width,height) #Add text box for title
        tf = txBox.text_frame
        page_count = 0
        for filename in sorted(os.listdir(os.path.join(pdf_path, folder))):# Access files within the folder, and sort them to appear in sam order on slide
          if filename.endswith(".pdf"):
  
            file = os.path.join(pdf_path,folder, filename) #pdf file path
            
            tf.text = folder
            print(folder)
            
            doc = fitz.open(file)
            for page in doc:  # iterate through the pages
              pix = page.get_pixmap(dpi=150)  # render page to an image
              pix.save(f"{pdf_path}/{folder}-{filename}.png")  # store image as a PNG
              
              
              slide.shapes.add_picture(f"{pdf_path}/{folder}-{filename}.png",left=Inches(page_count*2.5),top=Inches(1), width=Inches(2.5))
              print(page_count)
          page_count+=1
          
  
  prs.save(output_ppt_path)
  print(f"Successfully converted PDFs to PowerPoint presentation: {output_ppt_path}")
  return

#This function plots efel efeatures as bar graphs (in this case HH and HMM)
def plot_efeatures_bar(plot_folder,pfx):
    x = ['HH','HMM']
    data1 = np.genfromtxt('na12_orig1_efel.csv',delimiter=',')
    data2 = np.genfromtxt('na12_HMM_TF100923_efel.csv',delimiter=',')
    for i in range(1,15):
      y = [data1[1,i],data2[1,i]]
      print(y)
      fig, ax = plt.subplots()
      ax.bar(x,y,width=0.3,edgecolor='white')
      file_path_to_save=f'{plot_folder}{pfx}_{i}.pdf'
      plt.savefig(file_path_to_save, format='pdf')

    return


#This function takes efel features csvs which only have a header and one line of data and combines them into a single csv
def combine_efel_csvs(folder_path, output_file):
    combined_df = pd.DataFrame()
    
    for filename in os.listdir(folder_path):
        if filename.endswith(".csv"):
            file_path = os.path.join(folder_path, filename)
            df = pd.read_csv(file_path)
            combined_df = pd.concat([combined_df, df], ignore_index=True)
    
    combined_df.to_csv(output_file, index=False)

##Takes params text file and allows you to change the values for scanning etc...
def modify_dict_file(filename, changes):
  """
  Modifies values in a dictionary stored in a text file.

  Args:
      filename: The name of the text file containing the dictionary.
      changes: A dictionary containing key-value pairs where the key is the key to modify in the original dictionary and the value is the new value.

  Raises:
      ValueError: If the file cannot be opened or the content is not valid JSON.
  """

  try:
    # Open the file and read its content
    with open(filename, "r") as file:
      content = file.read()

    # Try to load the content as a dictionary
    try:
      data = eval(content)  # Assuming the file contains valid dictionary syntax
    except (NameError, SyntaxError):
      raise ValueError("Invalid dictionary format in the file.")

    # Modify values based on the provided changes dictionary
    for key, value in changes.items():
      if key not in data:
        print(f"Warning: Key '{key}' not found in the dictionary, skipping.")
      else:
        data[key] = value

    # Write the modified dictionary back to the file
    # with open(filename, "w") as file:
    #   file.write(repr(data))
    with open(filename, "w") as file:
      file.write(json.dumps(data, indent=2))  # Add indentation for readability (optional)

  except IOError as e:
    raise ValueError(f"Error opening or writing file: {e}")
##Using Modify_dict_file, setting args
filename = "./params/na16HH_TF2.txt"
changes = {
          "sh":8,
          "gbar":0.1,
          "tha":-59, #don't change
          "qa":4.5, #don't change
          "Ra":0.4,
          "Rb":0.124,
          "thi1":-80, #don't change
          "thi2":-80, #don't change
          "qd":5.4, #don't change
          "qg":5.4, #don't change
          "hmin":0.01,
          "mmin":0.02,
          "q10":2,
          "Rg":0.01,
          "Rd":0.03,
          "thinf":-80, #don't change
          "qinf":5.4, #don't change
          "vhalfs":-60, #don't change
          "a0s":0.0003,
          "zetas":12,
          "gms":0.2,
          "smax":10,
          "vvh":-58,
          "vvs":2,
          "ar2":1,
          "ena":55
          }

# modify_dict_file(filename, changes)

def combine_dictionaries(folder_path, new_file):
  """
  Combines dictionaries from .txt files in a folder into a single dictionary.

  Args:
      folder_path (str): Path to the folder containing .txt files with dictionaries.

  Returns:
      dict: A dictionary where file names are keys and dictionary contents are values.
  """
  combined_dict = {}
  for filename in os.listdir(folder_path):
    # Check for .txt files
    if filename.endswith(".txt"):
      file_path = os.path.join(folder_path, filename)
      try:
        # Open file and read content
        with open(file_path, 'r') as f:
          data = json.load(f)
        # Add dictionary to combined dict with filename as key
        combined_dict[filename] = data
      except (FileNotFoundError, json.JSONDecodeError) as e:
        print(f"Error processing file {filename}: {e}")
  with open (new_file ,'w') as file:
    # json.dump(combined_dict, file, separators=(',', ':')) #indent=4 if you want more vertical orientation
    # json.dump(combined_dict, file, indent=4) #indent=4 if you want more vertical orientation
    for key, value in combined_dict.items():
    # Write key-value pair on a separate line
      file.write(f'"{key}": {json.dumps(value)},\n')
  return combined_dict


## Plot states for 8 state HMM model
def plot_8states(csv_name,outfile_sfx,start=6500,stop=8500, ap_t=None, vm_t=None):
  df=pd.read_csv(csv_name, index_col=False)
  # Ensure start and end indices are within valid range
  if start < 0 or start >= len(df):
    raise ValueError("Invalid start index. Must be non-negative and less than data length.")
  if stop <= start or stop > len(df):
    raise ValueError("Invalid end index. Must be greater than start index and within data length.")
  
  dfsub = df.iloc[start:stop]

  fig1, ax1 = plt.subplots()
  ax1.plot(dfsub['c1'],label='c1', color='red', linewidth=0.7)
  ax1.plot(dfsub['c2'],label='c2', color='orange', linewidth=0.7)
  ax1.plot(dfsub['c3'],label='c3', color='pink', linewidth=0.7)
  ax1.plot(dfsub['i1'],label='i1', color='green', linewidth=0.7)
  ax1.plot(dfsub['i2'],label='i2', color='blue', linewidth=0.7)
  ax1.plot(dfsub['i3'],label='i3', color='cyan', linewidth=0.7)
  ax1.plot(dfsub['i4'],label='i4', color='purple', linewidth=0.7)
  ax1.plot(dfsub['o'],label='o', color='black', linewidth=0.7)

  ax2 = ax1.twinx()

  if ap_t is not None and vm_t is not None:
    y_min = min(np.min(vm_t[start:stop]), np.min(dfsub.min(axis=1)))
    y_max = max(np.max(vm_t[start:stop]), np.max(dfsub.max(axis=1)))
    ax2.set_ylim(bottom=y_min, top=y_max)
    ax2.plot(ap_t[start:stop],vm_t[start:stop], label='Vm', color='black',linewidth=2,linestyle='dashed') ##TF031424 changed linewidth
    ax2.set_ylabel('Vm', color='black')
  
  ax1.set_xlabel('timestep')
  ax1.set_ylabel('stateval')
  ax1.legend(loc='upper left')
  
  if ap_t is not None and vm_t is not None:
    ax2.legend(loc='upper right')
  
  # plt.legend()
  # plt.xticks(np.arange(0,len(df),1000), rotation=270)
  plt.xticks(np.arange(start,stop,500), rotation=270)
  plt.xlabel('timestep')
  # plt.ylabel('stateval')
  plt.title("States")
  # plt.savefig("8States"+".png", dpi=400)
  plt.savefig(f"/global/homes/t/tfenton/Neuron_general-2/Plots/Channel_state_plots/{start}-{stop}_{outfile_sfx}.png", dpi=400)



# This function takes a folder of efel data csvs and plots them in a heatmap. Each csv will be a new x-axis column in the heatmap
def efel_heatmaps(folder_path, output_folder):
  # rcParams['font.family'] = 'Arial'

  heatmap_data = {
    'dvdt Peak1 Height': pd.DataFrame(),
    'dvdt Peak2 Height': pd.DataFrame(),
    'dvdt Threshold': pd.DataFrame()
  }
  
  # sort filenames in alphabetical order
  filenames = sorted([f for f in os.listdir(folder_path) if f.endswith(".csv")])

  # Read each CSV file and extract the required columns
  for filename in filenames:
      file_path = os.path.join(folder_path, filename)
      df = pd.read_csv(file_path)
      key = os.path.splitext(filename)[0]  # Use filename without extension as the key
      
      heatmap_data['dvdt Peak1 Height'][key] = df.set_index('Type')['dvdt Peak1 Height']
      heatmap_data['dvdt Peak2 Height'][key] = df.set_index('Type')['dvdt Peak2 Height']
      heatmap_data['dvdt Threshold'][key] = df.set_index('Type')['dvdt Threshold']
  
  # Generate heatmaps for each feature
  for feature, data in heatmap_data.items():
    plt.figure(figsize=(10, 8))

    # Get min max values for each heatmap so each can have its own custom scale.
    min_val = data.min().min()
    max_val = data.max().max()
    
    cmap = mcolors.LinearSegmentedColormap.from_list("", ["blue", "white", "red"])

    x_labels = ['Left 2', 'Left 1', 'WT Baseline', 'Right 1', 'Right 2','Right 3', 'Right 4', 'Right 5']
    sns.heatmap(data, cmap=cmap, center=(min_val + max_val) / 2, annot=True, fmt=".2f", vmin=min_val, vmax=max_val)
    plt.title(f'{feature}')
    plt.xlabel('AIS Crossover Point Shift')
    plt.ylabel('% Nav1.2 : % Nav1.6 Ratio')
    plt.xticks(ticks=np.arange(len(x_labels)) + 0.5, labels=x_labels, rotation=90,fontsize=12, fontweight='bold')
    plt.yticks(rotation=0,fontsize=12, fontweight='bold')
    plt.tight_layout()
    
    # Save the heatmap
    output_file = os.path.join(output_folder, f'{feature}_heatmap.pdf')
    plt.savefig(output_file)
    plt.close()


# This function renames files and folders in a given directory by replacing ':' with '%'
def rename_files_and_folders(root_folder):
    for dirpath, dirnames, filenames in os.walk(root_folder):
        # Rename directories
        for dirname in dirnames:
            if ':' in dirname:
                new_dirname = dirname.replace(':', '%')
                os.rename(os.path.join(dirpath, dirname), os.path.join(dirpath, new_dirname))
        
        # Rename files
        for filename in filenames:
            if ':' in filename:
                new_filename = filename.replace(':', '%')
                os.rename(os.path.join(dirpath, filename), os.path.join(dirpath, new_filename))


def combine_pdfs_side_by_side(folder1, folder2, output_folder, output_suffix="_combined"):
    """
    Combines two sets of PDFs from two folders, placing corresponding pages side-by-side
    in a new PDF.  Only combines files with matching names.

    Args:
        folder1 (str): Path to the first folder containing PDFs.
        folder2 (str): Path to the second folder containing PDFs.
        output_folder (str): Path to the folder to save the combined PDF.
        output_suffix (str): Suffix to add to the output filename (default: "_combined").
    """

    os.makedirs(output_folder, exist_ok=True)  # Create output folder if it doesn't exist

    # Get the list of PDF files in folder1
    pdf_files1 = set(f for f in os.listdir(folder1) if f.endswith(".pdf"))

    # Get the list of PDF files in folder2
    pdf_files2 = set(f for f in os.listdir(folder2) if f.endswith(".pdf"))

    # Find the intersection of the two sets to get matching filenames
    matching_files = pdf_files1.intersection(pdf_files2)
    matching_files = sorted(list(matching_files))  # Sort for consistent order

    for filename in matching_files:
        pdf1_path = os.path.join(folder1, filename)
        pdf2_path = os.path.join(folder2, filename)

        # Create the output filename
        base_filename, ext = os.path.splitext(filename)
        output_pdf = os.path.join(output_folder, f"{base_filename}{output_suffix}.pdf")

        output_doc = fitz.open()  # Create a new PDF document

        try:
            doc1 = fitz.open(pdf1_path)
            doc2 = fitz.open(pdf2_path)

            # Ensure both PDFs have the same number of pages
            if doc1.page_count != doc2.page_count:
                print(f"Warning: {filename} has different page counts in {folder1} and {folder2}. Skipping.")
                doc1.close()
                doc2.close()
                continue

            # Determine the maximum height among all pages
            max_height = 0
            for page_num in range(doc1.page_count):
                page1 = doc1.load_page(page_num)
                page2 = doc2.load_page(page_num)
                max_height = max(max_height, page1.rect.height, page2.rect.height)

            for page_num in range(doc1.page_count):
                page1 = doc1.load_page(page_num)
                page2 = doc2.load_page(page_num)

                # Create a new page in the output document with the calculated height
                new_page = output_doc.new_page(width=page1.rect.width + page2.rect.width, height=max_height)

                # Calculate the width for each PDF page
                width1 = new_page.rect.width / 2
                width2 = new_page.rect.width / 2

                # Define rectangles for each PDF page
                rect1 = fitz.Rect(0, 0, width1, new_page.rect.height)
                rect2 = fitz.Rect(width1, 0, new_page.rect.width, new_page.rect.height)

                # Insert the PDF pages into the new page
                new_page.show_pdf_page(rect1, doc1, page_num)
                new_page.show_pdf_page(rect2, doc2, page_num)

            doc1.close()
            doc2.close()

        except Exception as e:
            print(f"Error processing {filename}: {e}")

        output_doc.save(output_pdf)
        output_doc.close()
        print(f"Successfully combined PDFs into {output_pdf}")


def combine_pdfs_flexible(folder, output_folder, match_length=None, match_suffix=None, output_suffix="_combined"):
    """
    Combines PDFs from a folder, placing pages side-by-side in a new PDF.
    Combines files with matching prefixes and suffixes of a specified length.

    Args:
        folder (str): Path to the folder containing PDFs.
        output_folder (str): Path to the folder to save the combined PDF.
        match_length (int, optional): Number of characters from the beginning of the filename to use for matching.
            If None, the entire filename (excluding extension) is used for matching. Defaults to None.
        match_suffix (str, optional): Characters at the end of the filename to use for matching. Defaults to None.
        output_suffix (str): Suffix to add to the output filename (default: "_combined").
    """

    os.makedirs(output_folder, exist_ok=True)  # Create output folder if it doesn't exist

    pdf_files = [f for f in os.listdir(folder) if f.endswith(".pdf")]
    
    # Group files based on the matching prefix and suffix
    grouped_files = {}
    for filename in pdf_files:
        base_filename, ext = os.path.splitext(filename)
        prefix = base_filename[:match_length] if match_length else base_filename
        suffix = base_filename[-len(match_suffix):] if match_suffix else ""  # Extract suffix

        # Combine prefix and suffix for grouping
        match_key = (prefix, suffix)

        if match_suffix:
          if len(base_filename) < len(match_suffix):
            continue #skip file if it's shorter than the suffix

          if base_filename[-len(match_suffix):] != match_suffix:
            continue #skip file if suffix doesn't match

        if match_key not in grouped_files:
            grouped_files[match_key] = []
        grouped_files[match_key].append(os.path.join(folder, filename))

    for (match_prefix, match_suffix), file_paths in grouped_files.items():
        if len(file_paths) < 2:
            print(f"Skipping {match_prefix}_{match_suffix}: Not enough files to combine.")
            continue

        # Sort file paths based on the numerical value in "LVA###"
        def extract_number(file_path):
          match = re.search(r"HVA+SKE2-(\d+(?:\.\d+)?)", file_path)  # Match integers or floats
          if match:
              try:
                  num_str = match.group(1)
                  if '.' in num_str:
                      return float(num_str)  # Handle floats
                  else:
                      return int(num_str)  # Handle integers
              except ValueError:
                  return float('inf')
          return float('inf')  # If "LVA###" not found, put it at the end

        file_paths.sort(key=extract_number)

        # Create the output filename
        output_pdf = os.path.join(output_folder, f"{match_prefix}{output_suffix}.pdf")
        output_doc = fitz.open()  # Create a new PDF document

        try:
            # Open all documents
            docs = [fitz.open(pdf_path) for pdf_path in file_paths]
            filenames = [os.path.basename(pdf_path) for pdf_path in file_paths] #get filenames

            # Check if all PDFs have the same number of pages
            num_pages = docs[0].page_count
            if not all(doc.page_count == num_pages for doc in docs):
                print(f"Warning: {match_prefix} has PDFs with different page counts. Skipping.")
                for doc in docs:
                    doc.close()
                continue

            # Determine the maximum height among all pages in all documents
            max_height = 0
            for page_num in range(num_pages):
                for doc in docs:
                    page = doc.load_page(page_num)
                    max_height = max(max_height, page.rect.height)

            # Combine pages side by side
            for page_num in range(num_pages):
                # Calculate total width
                total_width = sum(doc.load_page(page_num).rect.width for doc in docs)
                new_page = output_doc.new_page(width=total_width, height=max_height+30) #add space for text
                
                # Place each page side by side
                current_x = 0
                for i, doc in enumerate(docs):
                    page = doc.load_page(page_num)
                    width = page.rect.width
                    rect = fitz.Rect(current_x, 30, current_x + width, new_page.rect.height) #Shift down 20 to make room for text
                    new_page.show_pdf_page(rect, doc, page_num)

                     # Add filename above the image
                    text_rect = fitz.Rect(current_x+10, 20, current_x + width, 30)
                    if len(filenames[i]) >= 12:
                        new_page.insert_text(text_rect.tl, filenames[i][7:21], color=(0, 0, 0), fontsize=8)  #black text, fontsize 8
                    else:
                        new_page.insert_text(text_rect.tl, filenames[i], color=(0, 0, 0), fontsize=8)  # Use the full name if it's too short

                    current_x += width

            # Close all documents
            for doc in docs:
                doc.close()

        except Exception as e:
            print(f"Error processing {match_prefix}: {e}")

        output_doc.save(output_pdf)
        output_doc.close()
        print(f"Successfully combined PDFs into {output_pdf}")


def plot_and_auc_for_folder(folder, save_plots=True, plot_dir_name='plots_auc',csv_out='FI_AUCs'):
    """
    For each txt file in the folder, plot the curve and calculate area under the curve.

    Args:
        folder (str): Path to folder containing txt files (each with 140 y-values).
        save_plots (bool): If True, saves each plot as PNG in plot_dir.
        plot_dir (str): Directory to save plots.

    Returns:
        dict: {filename: auc_value}
    """
    plot_output = os.path.join(folder, plot_dir_name)
    if save_plots and not os.path.exists(plot_output):
        os.makedirs(plot_output)

    auc_dict = {}
    x = np.linspace(-0.4, 1, 140)

    for fname in os.listdir(folder):
        if fname.endswith('.txt'):
            fpath = os.path.join(folder, fname)
            with open(fpath, 'r') as f:
                # Try to read as list of floats
                y = eval(f.read())
                y = np.array(y, dtype=float)
                if y.shape[0] != 140:
                    print(f"Warning: {fname} does not have 140 values, skipping.")
                    continue

            auc = np.trapz(y, x)
            auc_dict[fname] = auc
            # print(auc_dict)
            plt.figure()
            plt.plot(x, y, label=fname)
            plt.title(f"{fname} (AUC={auc:.2f})")
            plt.xlabel("Injected Current (nA)")
            plt.ylabel("APs per 500ms epoch")
            plt.legend()
            if save_plots:
                plt.savefig(os.path.join(plot_output, f"{fname}.png"))
            plt.close()
    
    # Write AUC results to CSV if requested
    if csv_out:
        with open(f'{plot_output}/{csv_out}.csv', 'w', newline='') as csvfile:
            writer = csv.writer(csvfile)
            writer.writerow(['filename', 'auc'])
            for fname, auc in auc_dict.items():
                writer.writerow([fname, auc])

    return auc_dict


def first_positive_time_in_csvs(input_folder, output_csv):

  '''
  #################################################################################################
  ##### Example usage to first get many csvs (1 for each segment, in this case for section axon[0])
  ##################################################################################################
  simwt = tf.Na12Model_TF(ais_nav12_fac=5.76,ais_nav16_fac=1.08*0.6,
                            nav12=1.1*1.1,nav16=1.43*1.2, somaK=0.022, KP=3.9375, KT=5,
                            ais_ca = 43*0.5,ais_Kca = 0.25,
                            soma_na16=0.8,soma_na12=2.56,node_na = 1,
                            dend_nav12=1,
                            na12name = 'na12annaTFHH2',mut_name = 'na12annaTFHH2',na12mechs = ['na12','na12mut'],
                            na16name = 'na16HH_TF2',na16mut_name = 'na16HH_TF2',na16mechs=['na16','na16mut'],params_folder = './params/',
                            plots_folder = f'{root_path_out}/{path}', update=True, fac=None)



for i in range(1,121):
  # seg = i/121
  seg = i/121

  sim_config = {
            'section': 'axon',
            'section_num': 0,
            'segment': seg,
            'currents'  : ['ica_Ca_HVA','ica_Ca_LVAst','ik_K_Pst','ik_K_Tst','ik_SK_E2','ik_SKv3_1',
                               'na12.ina_ina','na12mut.ina_ina','na16.ina_ina','na16mut.ina_ina','i_pas'], #AIS (no Ih)
            'current_names' : ['Ca_HVA','Ca_LVAst','K_Pst','K_Tst','SK_E2','SKv3_1','Na12','Na12 MUT','Na16 WT','Na16 MUT','pas'],
            'ionic_concentrations' :["ki", "nai","cai"]
            }
  Vm_array, I, t, stim = simwt.get_stim_raw_data(stim_amp=0.5, dt=0.005, stim_dur=120, sim_config=sim_config)
  # Inside your for loop, for each segment:
  df = pd.DataFrame({'time': t, 'Vm': Vm_array})
  csv_path = f"{root_path_out}/Vm_segment_{seg}.csv"
  df.to_csv(csv_path, index=False)
  print(f"Saved Vm for segment {seg} to {csv_path}")
'''

  # input_folder = "./Plots/12HH16HH/28-AP_initiation"  # Change to your folder path
  output_csv = f"{input_folder}/{output_csv}"

  results = []

  for fname in os.listdir(input_folder):
      if fname.endswith(".csv"):
          fpath = os.path.join(input_folder, fname)
          df = pd.read_csv(fpath)
          # Assumes columns: 'time', 'Vm'
          idx = (df['Vm'] > -40).idxmax()
          if df['Vm'][idx] > -40:
              results.append([fname, df['time'][idx], df['Vm'][idx]])
          else:
              results.append([fname, None, None])

  out_df = pd.DataFrame(results, columns=['filename', 'first_positive_time', 'voltage'])
  out_df.to_csv(output_csv, index=False)
  print(f"Saved summary to {output_csv}")


# combined_dict = combine_dictionaries(folder_path='/global/homes/t/tfenton/Neuron_general-2/params/na16_HOF_params_JSON', new_file='/global/homes/t/tfenton/Neuron_general-2/params/na16_HOF_params_JSON/combined3.json')


# plot_8states(csv_name="pandas_states.csv", outfile_sfx="8st_062824")



#combine_pdfs(folder_path='/global/homes/t/tfenton/Neuron_general-2/Plots/12HMM16HH_TF/AllSynthMuts_121223/', out_sfx='traceDVDTfi_121323')


#combine_and_sort_ef_csvs(root_path='/global/homes/t/tfenton/Neuron_general-2/Plots/12HMM16HH_TF/AllSynthMuts_121223/', out_sfx='allsynthmuts_121323_sorted')

# make_ppt_from_pdf2(pdf_path='./Plots/12HH16HH/5-newAIS_raiseDVDT/37-AdilHHvariants_800sweep_091824',
#                   output_ppt_path='./Plots/12HH16HH/5-newAIS_raiseDVDT/37-AdilHHvariants_800sweep_091824/HHmuts.pptx')

#plot_efeatures_bar(plot_folder='/global/homes/t/tfenton/Neuron_general-2/Plots/12HMM16HH_TF/ManuscriptFigs/efeatures',pfx='soma')


# combine_efel_csvs(folder_path='./Plots/12HH16HH/23-PaperPlots/3-SynthMuts_052325', output_file='./Plots/12HH16HH/23-PaperPlots/3-SynthMuts_052325/synthmuts_combined_efel_052725.csv')
# efel_heatmaps('./Plots/12HH16HH/10-KevinRtR_chandensities/11-ShiftAIS/30-newCombinedCsvs/Updated_EFEL_peak2', './Plots/12HH16HH/10-KevinRtR_chandensities/11-ShiftAIS/30-newCombinedCsvs/Updated_EFEL_peak2')
# rename_files_and_folders('./Plots/12HH16HH/10-KevinRtR_chandensities/11-ShiftAIS')

# combine_pdfs_side_by_side(folder1='./Plots/12HH16HH/6-October2024Model/1-AdilHHvariants_800sweep_101524', folder2='./Plots/12HH16HH/15-AdilMuts_newModel_042425', output_folder='./Plots/12HH16HH/15-AdilMuts_newModel_042425/Combined/')
# combine_pdfs_flexible(folder='./Plots/12HH16HH/18-SynthMuts_newmodel_shortstim/5-HVA+SKE2', 
#                       output_folder='./Plots/12HH16HH/18-SynthMuts_newmodel_shortstim/5-HVA+SKE2/Combined', 
#                       match_length=7,
#                       match_suffix='wtvmut')


# plot_and_auc_for_folder(folder='./Plots/12HH16HH/23-PaperPlots/Fig4',
#                         save_plots=True)
first_positive_time_in_csvs(input_folder='./Plots/12HH16HH/29-AP_initiation_higher_res', output_csv=f'./Plots/12HH16HH/29-AP_initiation_higher_res/first_time_over-40.csv')